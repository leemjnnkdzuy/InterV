import pptx
import json

prs = pptx.Presentation(r'd:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v18.pptx')

slides_to_inspect = [8, 11, 12, 13, 14, 15, 16, 22, 23, 24, 25, 26] # 0-indexed: slides 9, 12, 13, 14, 15, 16, 17, 23, 24, 25, 26, 27
out_data = {}

for s_idx in slides_to_inspect:
    slide = prs.slides[s_idx]
    slide_num = s_idx + 1
    out_data[f"Slide_{slide_num}"] = []
    for sh_idx, shape in enumerate(slide.shapes):
        if shape.has_text_frame:
            sh_entry = {
                "shape_idx": sh_idx,
                "shape_name": shape.name,
                "paragraphs": []
            }
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                p_entry = {
                    "p_idx": p_idx,
                    "text": p.text,
                    "runs": [{"r_idx": r_idx, "text": r.text} for r_idx, r in enumerate(p.runs)]
                }
                sh_entry["paragraphs"].append(p_entry)
            out_data[f"Slide_{slide_num}"].append(sh_entry)

with open(r'd:\project\InterV\Review\key_slides_runs.json', 'w', encoding='utf-8') as f:
    json.dump(out_data, f, ensure_ascii=False, indent=2)

print("Saved key_slides_runs.json successfully")
