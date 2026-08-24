# -*- coding: utf-8 -*-
"""
Script to update Slide 15 in PPTX and corresponding Word documentation (Kich Ban Thuyet Trinh & KLTN Thesis)
to accurately reflect the two-phase Multimodal Speech Pipeline in InterV:
1. Phase 1 (Realtime Turn-taking): Microphone -> AssemblyAI Streaming STT -> Local Faster-Whisper Fallback -> Prompt/Audio Context to AI Provider (DeepSeek) for seamless Lookahead Question Generation and TTS without conversational gaps.
2. Phase 2 (Post-Interview Sequential Processing): Saved Audio Chunks -> Sequential Faster-Whisper (high-accuracy ASR) -> Prosody/Signal Processing Algorithms (pacing, WPM, pause gaps) -> SenseVoice LID/SER/AED (Observation-only) -> Comprehensive Delivery/Confidence/Fluency Coaching Metrics.
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import docx
import json

def update_pptx(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
    
    prs = pptx.Presentation(pptx_path)
    slide = prs.slides[14] # Slide 15 (0-indexed 14)
    
    # Shape 0: Title
    shape0 = slide.shapes[0]
    shape0.text_frame.clear()
    p0 = shape0.text_frame.paragraphs[0]
    p0.text = "Pipeline tiếng nói nhiều tầng: Realtime & Hậu kỳ phỏng vấn"
    p0.font.name = "Cambria"
    p0.font.size = Pt(21)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    # Shape 2: Subtitle
    shape2 = slide.shapes[2]
    shape2.text_frame.clear()
    p2 = shape2.text_frame.paragraphs[0]
    p2.text = "AssemblyAI streaming • Faster-Whisper fallback & batch • SenseVoice LID/SER/AED"
    p2.font.name = "Montserrat"
    p2.font.size = Pt(9.75)
    p2.font.bold = False
    p2.font.color.rgb = RGBColor(0x52, 0x60, 0x52)
    
    # Shape 3: Header Row 1
    shape3 = slide.shapes[3]
    shape3.text_frame.clear()
    p3 = shape3.text_frame.paragraphs[0]
    p3.text = "1. LUỒNG REALTIME TRONG PHỎNG VẤN (TƯƠNG TÁC LIỀN MẠCH)"
    p3.font.name = "Montserrat"
    p3.font.size = Pt(8.25)
    p3.font.bold = True
    p3.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
    
    # Shape 4: Clear / Hide unused badge
    shape4 = slide.shapes[4]
    shape4.text_frame.clear()
    
    # Shape 5: Header Row 2
    shape5 = slide.shapes[5]
    shape5.text_frame.clear()
    p5 = shape5.text_frame.paragraphs[0]
    p5.text = "2. LUỒNG XỬ LÝ HẬU KỲ SAU PHỎNG VẤN (ĐÁNH GIÁ CHUYÊN SÂU)"
    p5.font.name = "Montserrat"
    p5.font.size = Pt(8.25)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
    
    # Shape 6 (Row 1 Box 1): Microphone
    shape6 = slide.shapes[6]
    shape6.text_frame.clear()
    p6_1 = shape6.text_frame.paragraphs[0]
    p6_1.text = "Microphone"
    p6_1.font.name = "Montserrat"
    p6_1.font.size = Pt(9.5)
    p6_1.font.bold = True
    p6_1.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    p6_2 = shape6.text_frame.add_paragraph()
    p6_2.text = "Thu âm AudioWorklet (16kHz PCM)"
    p6_2.font.name = "Montserrat"
    p6_2.font.size = Pt(8.25)
    p6_2.font.color.rgb = RGBColor(0x52, 0x60, 0x52)
    
    # Shape 7 (Row 1 Box 2): AssemblyAI Streaming + Fallback
    shape7 = slide.shapes[7]
    shape7.text_frame.clear()
    p7_1 = shape7.text_frame.paragraphs[0]
    p7_1.text = "AssemblyAI Streaming"
    p7_1.font.name = "Montserrat"
    p7_1.font.size = Pt(9.5)
    p7_1.font.bold = True
    p7_1.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    p7_2 = shape7.text_frame.add_paragraph()
    p7_2.text = "STT thời gian thực (WebSocket)\n[Lỗi → Faster-Whisper fallback]"
    p7_2.font.name = "Montserrat"
    p7_2.font.size = Pt(8.0)
    p7_2.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
    
    # Shape 8 (Row 1 Box 3): AI Provider
    shape8 = slide.shapes[8]
    shape8.text_frame.clear()
    p8_1 = shape8.text_frame.paragraphs[0]
    p8_1.text = "AI Provider (DeepSeek)"
    p8_1.font.name = "Montserrat"
    p8_1.font.size = Pt(9.5)
    p8_1.font.bold = True
    p8_1.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    p8_2 = shape8.text_frame.add_paragraph()
    p8_2.text = "Có text tức thời → Sinh câu hỏi thích ứng & TTS mượt mà, không rời rạc"
    p8_2.font.name = "Montserrat"
    p8_2.font.size = Pt(8.0)
    p8_2.font.color.rgb = RGBColor(0x52, 0x60, 0x52)
    
    # Shape 9 (Row 2 Box 1): Faster-Whisper batch + Prosody Algorithms
    shape9 = slide.shapes[9]
    shape9.text_frame.clear()
    p9_1 = shape9.text_frame.paragraphs[0]
    p9_1.text = "Faster-Whisper tuần tự"
    p9_1.font.name = "Montserrat"
    p9_1.font.size = Pt(9.5)
    p9_1.font.bold = True
    p9_1.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    p9_2 = shape9.text_frame.add_paragraph()
    p9_2.text = "STT chính xác cao + Thuật toán: nhịp độ, tốc độ (WPM), khoảng lặng"
    p9_2.font.name = "Montserrat"
    p9_2.font.size = Pt(8.0)
    p9_2.font.color.rgb = RGBColor(0x52, 0x60, 0x52)
    
    # Shape 10 (Row 2 Box 2): SenseVoice LID / SER / AED
    shape10 = slide.shapes[10]
    shape10.text_frame.clear()
    p10_1 = shape10.text_frame.paragraphs[0]
    p10_1.text = "SenseVoice (Observation only)"
    p10_1.font.name = "Montserrat"
    p10_1.font.size = Pt(9.0)
    p10_1.font.bold = True
    p10_1.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    p10_2 = shape10.text_frame.add_paragraph()
    p10_2.text = "LID (Ngôn ngữ) • SER (Cảm xúc)\nAED (Tiếng thở, cười, filler words)"
    p10_2.font.name = "Montserrat"
    p10_2.font.size = Pt(8.0)
    p10_2.font.color.rgb = RGBColor(0x52, 0x60, 0x52)
    
    # Shape 11 (Row 2 Box 3): Evaluation & Coaching
    shape11 = slide.shapes[11]
    shape11.text_frame.clear()
    p11_1 = shape11.text_frame.paragraphs[0]
    p11_1.text = "Đánh giá phong thái"
    p11_1.font.name = "Montserrat"
    p11_1.font.size = Pt(9.5)
    p11_1.font.bold = True
    p11_1.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
    
    p11_2 = shape11.text_frame.add_paragraph()
    p11_2.text = "Mức độ tự tin, lưu loát, nhịp điệu diễn đạt hỗ trợ coaching"
    p11_2.font.name = "Montserrat"
    p11_2.font.size = Pt(8.0)
    p11_2.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
    
    # Adjust connector 14 to connect Box 4 (Shape 9) to Box 5 (Shape 10) horizontally
    conn14 = slide.shapes[14]
    conn14.left = 2809875
    conn14.top = 3799691
    conn14.width = 857250
    conn14.height = 0
    
    # Hide connector 15 (old vertical line)
    conn15 = slide.shapes[15]
    conn15.width = 0
    conn15.height = 0
    
    # Update Slide Notes
    speaking_note = (
        "[08:00 - 09:30] SLIDE 15: PIPELINE TIẾNG NÓI ĐA TẦNG (REALTIME & HẬU KỲ)\n\n"
        "Lời nói:\n"
        "\"Kính thưa Thầy Cô, pipeline tiếng nói của InterV được thiết kế chuyên biệt theo 2 pha rõ rệt:\n\n"
        "1. PHA 1: LUỒNG REALTIME TRONG PHỎNG VẤN (TƯƠNG TÁC MƯỢT MÀ, KHÔNG RỜI RẠC)\n"
        "• Luồng âm thanh từ microphone được AudioWorklet thu và truyền trực tiếp qua WebSocket tới AssemblyAI để thực hiện STT theo thời gian thực.\n"
        "• Cơ chế Fallback cục bộ: Khi gặp sự cố mạng hoặc lỗi API, hệ thống tự động fallback về model Faster-Whisper chạy trực tiếp trên Python backend.\n"
        "• Mục đích: Luôn bảo đảm có dữ liệu transcript tức thời cho AI Provider (DeepSeek) tiếp tục hiểu ngữ cảnh, nói tiếp qua TTS và sinh các câu hỏi thích ứng (Lookahead Engine) liền mạch, không làm câu hỏi bị rời rạc hay ngắt quãng phiên phỏng vấn.\n\n"
        "2. PHA 2: LUỒNG XỬ LÝ HẬU KỲ SAU PHỎNG VẤN (ĐÁNH GIÁ CHUYÊN SÂU & COACHING)\n"
        "• Sau khi phỏng vấn xong, toàn bộ các đoạn audio gốc đã lưu sẽ được xử lý tuần tự (sequential batch) bằng Faster-Whisper để tái tạo STT chính xác cao nhất.\n"
        "• Hệ thống áp dụng các thuật toán xử lý tín hiệu âm thanh để đo lường khách quan nhịp độ, tốc độ nói (WPM), và các khoảng trống/khoảng lặng (pause duration) trong câu nói.\n"
        "• Tiếp theo, audio được đưa qua mô hình SenseVoice theo nguyên tắc 'Observation only' để bóc tách: LID (ngôn ngữ), SER (cảm xúc giọng nói), và AED (sự kiện âm thanh: tiếng cười, tiếng ho, thở dài, từ đệm/filler words).\n"
        "• Dữ liệu này được tổng hợp để đánh giá khách quan mức độ tự tin, tính lưu loát, và phong thái diễn đạt nhằm hỗ trợ coaching cho ứng viên mà không suy diễn tâm lý tùy tiện.\""
    )
    
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        slide.notes_slide.notes_text_frame.text = speaking_note
        
    prs.save(pptx_path)
    print(f"Successfully updated Slide 15 in {pptx_path}")

def update_kich_ban_docx():
    docx_path = r"d:\project\InterV\Review\Kich_Ban_Thuyet_Trinh_KLTN_InterV_15Phut.docx"
    if not os.path.exists(docx_path):
        print(f"File not found: {docx_path}")
        return
    
    doc = docx.Document(docx_path)
    
    # 1. Update paragraph 26 & 27
    for p in doc.paragraphs:
        txt = p.text.strip()
        if "▶ Slide 15: Pipeline Xử lý Giọng nói Đa tầng" in txt:
            p.text = "▶ Slide 15: Pipeline Xử lý Giọng nói Đa tầng (Realtime & Hậu kỳ phỏng vấn) [04:30 - 06:30 (2 phút - 120 giây)]"
        elif "Tiếp theo, em xin trình bày về Pipeline Giọng nói Đa tầng của InterV" in txt:
            p.text = (
                "Lời nói: \"Tiếp theo, em xin trình bày về Pipeline Giọng nói Đa tầng của InterV. "
                "Hệ thống được thiết kế chuyên biệt theo 2 pha rõ rệt để vừa đảm bảo độ trễ 0ms trong lúc đàm thoại, "
                "vừa đạt độ chính xác tối đa trong báo cáo đánh giá:\n\n"
                "1. Pha 1 - Luồng Realtime trong phỏng vấn (Tương tác mượt mà, không rời rạc):\n"
                "• Trình duyệt mở AudioWorklet thu âm PCM 16kHz và stream trực tiếp qua WebSocket tới AssemblyAI bằng temporary token ngắn hạn để STT thời gian thực, chữ chạy mượt mà trên màn hình.\n"
                "• Cơ chế Fallback cục bộ: Nếu gặp sự cố mạng hoặc lỗi kết nối, hệ thống tự động fallback về Faster-Whisper cục bộ trên Python backend. Điều này đảm bảo AI Provider (DeepSeek) luôn có dữ liệu câu trả lời tức thời để tiếp tục đối thoại, nói tiếp qua TTS và sinh câu hỏi thích ứng (Lookahead Engine) liền mạch, không làm câu hỏi bị rời rạc hay ngắt quãng trải nghiệm của ứng viên.\n\n"
                "2. Pha 2 - Luồng Xử lý Hậu kỳ sau phỏng vấn (Đánh giá chuyên sâu & Coaching):\n"
                "• Sau khi buổi phỏng vấn kết thúc, toàn bộ các đoạn audio đã lưu sẽ được xử lý tuần tự (sequential batch) bằng Faster-Whisper để tái tạo bản transcript với độ chính xác cao nhất kèm word-level timestamps.\n"
                "• Hệ thống áp dụng các thuật toán xử lý tín hiệu âm thanh để đo lường khách quan nhịp độ (pacing), tốc độ nói (Words Per Minute), và các khoảng trống/khoảng lặng (pause duration) trong câu nói.\n"
                "• Sau đó, audio được đưa qua mô hình SenseVoice theo nguyên tắc 'Observation only' để bóc tách: LID (ngôn ngữ), SER (cảm xúc giọng nói: tự tin, trung tính, căng thẳng), và AED (sự kiện âm thanh: tiếng thở dài, tiếng cười, tiếng ho, từ đệm/filler words).\n"
                "• Toàn bộ các chỉ số này được tổng hợp để đánh giá khách quan mức độ tự tin, tính lưu loát và phong thái diễn đạt, hỗ trợ ứng viên hoàn thiện kỹ năng phỏng vấn mà tuyệt đối không suy diễn tâm lý tùy tiện.\""
            )
            
    # 2. Update Table summaries
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if "Pipeline Giọng nói Đa tầng (AssemblyAI, Faster-Whisper, Vbee TTS)" in cell.text:
                    cell.text = cell.text.replace(
                        "Pipeline Giọng nói Đa tầng (AssemblyAI, Faster-Whisper, Vbee TTS)",
                        "Pipeline Giọng nói 2 Pha (Realtime STT + Fallback & Hậu kỳ Faster-Whisper + SenseVoice Coaching)"
                    )
                if "Nhấn mạnh cơ chế Fallback Faster-Whisper và Thư viện phiên âm thuật ngữ IT trong TTS." in cell.text:
                    cell.text = cell.text.replace(
                        "Nhấn mạnh cơ chế Fallback Faster-Whisper và Thư viện phiên âm thuật ngữ IT trong TTS.",
                        "Nhấn mạnh luồng 2 pha: Pha 1 Realtime STT + Fallback giúp AI Provider sinh câu hỏi liền mạch; Pha 2 Xử lý tuần tự Faster-Whisper + Thuật toán đo nhịp độ, khoảng lặng + SenseVoice Observation only cho Coaching tự tin, lưu loát."
                    )
                if "• Cấp temporary token ngắn hạn cho browser (bảo mật tuyệt đối)." in cell.text:
                    cell.text = (
                        "• Pha 1: Realtime STT qua AssemblyAI WS + Faster-Whisper fallback cục bộ giúp AI Provider sinh câu hỏi liền mạch.\n"
                        "• Pha 2: Xử lý tuần tự Faster-Whisper (STT chính xác) + Thuật toán đo nhịp độ, tốc độ WPM, khoảng lặng.\n"
                        "• SenseVoice Observation only (LID/SER/AED) đánh giá tự tin, lưu loát cho coaching."
                    )

    try:
        doc.save(docx_path)
        print(f"Successfully updated {docx_path}")
    except PermissionError:
        alt_path = r"d:\project\InterV\Review\Kich_Ban_Thuyet_Trinh_KLTN_InterV_15Phut_updated.docx"
        doc.save(alt_path)
        print(f"File was locked by Microsoft Word. Saved updated version to: {alt_path}")

def update_thesis_v14():
    docx_path = r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v14.docx"
    if not os.path.exists(docx_path):
        print(f"File not found: {docx_path}")
        return
    
    doc = docx.Document(docx_path)
    
    # 1. Update paragraphs
    for idx, p in enumerate(doc.paragraphs):
        txt = p.text
        
        # Front Matter Vietnamese Abstract (Paragraph 42)
        if "AssemblyAI và Faster-Whisper tạo lớp transcript dự phòng" in txt:
            p.text = txt.replace(
                "AssemblyAI và Faster-Whisper tạo lớp transcript dự phòng",
                "luồng âm thanh được phiên âm thời gian thực bằng AssemblyAI Streaming (với Faster-Whisper fallback cục bộ) giúp AI provider duy trì hội thoại và sinh câu hỏi thích ứng liền mạch; sau phỏng vấn, audio được xử lý tuần tự bằng Faster-Whisper để STT chính xác cao kết hợp thuật toán đo nhịp độ, tốc độ, khoảng lặng, sau đó chuyển tiếp qua SenseVoice (LID/SER/AED theo nguyên tắc observation-only) để nhận dạng trạng thái, biểu đạt và cảm xúc nhằm đánh giá mức độ tự tin, lưu loát hỗ trợ coaching"
            )
            
        # Front Matter English Abstract (Paragraph 48)
        if "AssemblyAI and Faster-Whisper provide transcript paths; SenseVoice produces observable delivery signals" in txt:
            p.text = txt.replace(
                "AssemblyAI and Faster-Whisper provide transcript paths; SenseVoice produces observable delivery signals without treating emotion tags as psychological or hiring decisions",
                "audio is transcribed in realtime via AssemblyAI Streaming with local Faster-Whisper fallback to feed the AI provider seamlessly for adaptive question generation without conversational gaps; post-interview audio is processed sequentially using Faster-Whisper for high-accuracy ASR and prosody algorithms (pacing, speech rate, pause gaps), followed by SenseVoice LID/SER/AED under the observation-only rule to evaluate confidence and fluency coaching metrics"
            )
            
        # Paragraph 87 (Overview)
        if "Một pipeline tiếng nói phân vai giữa AssemblyAI streaming, Faster-Whisper fallback" in txt:
            p.text = (
                "- Một pipeline tiếng nói 2 pha: Pha 1 phiên âm realtime bằng AssemblyAI streaming (với Faster-Whisper fallback cục bộ) "
                "đảm bảo dữ liệu tức thời cho AI provider sinh câu hỏi liên tục không bị rời rạc; "
                "Pha 2 xử lý tuần tự hậu kỳ bằng Faster-Whisper để STT chính xác cao, kết hợp thuật toán đo nhịp độ, tốc độ và khoảng lặng, "
                "sau đó qua SenseVoice (LID/SER/AED theo nguyên tắc observation-only) để đánh giá khách quan mức độ tự tin, lưu loát phục vụ coaching."
            )
            
        # Section 1.5 - Paragraph 192
        if "Một mô hình duy nhất không tối ưu cho mọi nhiệm vụ âm thanh. InterV phân tách transcript trực tuyến" in txt:
            p.text = (
                "Một mô hình duy nhất không tối ưu cho mọi nhiệm vụ âm thanh. InterV thiết kế pipeline tiếng nói chuyên biệt theo hai pha rõ rệt: "
                "(1) Pha thời gian thực (Interactive Realtime Turn): Luồng âm thanh từ microphone được AudioWorklet thu và truyền trực tiếp qua WebSocket tới AssemblyAI để nhận transcript tức thời. "
                "Khi xảy ra sự cố mạng hoặc lỗi API, hệ thống kích hoạt Faster-Whisper Fallback cục bộ trên Python backend để cung cấp dữ liệu văn bản cho AI Provider (DeepSeek), "
                "giúp AI tiếp tục nói tiếp qua TTS và sinh các câu hỏi thích ứng (Lookahead Adaptive Question) liền mạch, ngăn ngừa hoàn toàn tình trạng câu hỏi bị rời rạc hay đứt quãng phiên phỏng vấn. "
                "(2) Pha xử lý hậu kỳ (Post-Interview Sequential Processing): Sau khi phiên phỏng vấn hoàn tất, toàn bộ các đoạn audio đã lưu được xử lý tuần tự bằng Faster-Whisper "
                "với chất lượng cao nhất để tái tạo STT chính xác tuyệt đối. Hệ thống áp dụng các thuật toán xử lý tín hiệu âm thanh để đo lường khách quan nhịp độ (pacing), tốc độ nói (WPM), "
                "và các khoảng trống/khoảng lặng (pause duration) trong câu nói. Sau đó, audio được đưa qua mô hình SenseVoice (LID/SER/AED theo nguyên tắc Observation Only) "
                "để nhận dạng trạng thái, biểu đạt âm thanh (tiếng cười, tiếng ho, thở dài, từ đệm filler words) và cảm xúc giọng nói nhằm đánh giá mức độ tự tin, lưu loát phục vụ coaching cho ứng viên."
            )
            
        # Section 1.5 - Paragraph 200 & 201
        if "Tiền xử lý và tổng hợp nhiều đoạn: Âm thanh được nhận theo chunk" in txt:
            p.text = (
                "Tiền xử lý và tổng hợp nhiều đoạn trong xử lý hậu kỳ: Sau khi kết thúc phỏng vấn, các đoạn âm thanh đã lưu được đưa vào xử lý tuần tự qua Faster-Whisper "
                "để trích xuất transcript chính xác cao kèm word-level timestamps. Dựa trên các mốc thời gian này, thuật toán âm học phân tích cấu trúc nhịp điệu: "
                "tính toán tốc độ nói WPM, phát hiện các khoảng lặng vượt ngưỡng (pause detection) và phân loại các khoảng trống ngắt nghỉ. "
                "Tiếp đó, audio sample được chuyển qua SenseVoice để trích xuất LID/SER/AED ngoài event loop bằng asyncio.to_thread."
            )
            
        if "Baseline âm thanh sử dụng duration và số từ để ước lượng tốc độ nói" in txt:
            p.text = (
                "Đánh giá phong thái và mức độ tự tin, lưu loát: Thuật toán âm học kết hợp số từ, thời lượng phát âm và khoảng lặng để đo lường nhịp độ (pacing) và tính lưu loát (fluency). "
                "SenseVoice cung cấp các quan sát phong phú về sự kiện âm thanh (tiếng thở, cười, từ đệm) và cảm xúc bề mặt (SER) theo nguyên tắc Observation Only. "
                "Các chỉ số này được tổng hợp vào vocal_delivery và báo cáo coaching, giúp ứng viên nhận biết thói quen diễn đạt của mình mà không bị gán nhãn tâm lý võ đoán."
            )
            
        # Section 2.6 - Paragraph 401 & 402
        if "Trong lúc trả lời, trình duyệt thu âm và nhận token streaming từ CreateStreamingToken" in txt:
            p.text = (
                "Pipeline âm thanh được phân tách thành 2 pha vận hành rõ ràng: "
                "Trong lúc phỏng vấn (Pha Realtime), trình duyệt thu âm và mở kết nối WebSocket tới AssemblyAI thông qua token ngắn hạn từ CreateStreamingToken. "
                "Transcript trực tuyến hiển thị cho người dùng và cung cấp ngay lập tức cho AI Provider. Nếu xảy ra lỗi streaming hoặc mất kết nối, "
                "hệ thống tự động kích hoạt TranscribeAudio (Faster-Whisper fallback cục bộ) để AI Provider tiếp tục sinh câu hỏi thích ứng mà không làm gián đoạn hội thoại. "
                "Sau khi kết thúc phiên (Pha Hậu kỳ), toàn bộ audio được xử lý tuần tự bằng Faster-Whisper để STT chính xác cao, kết hợp các thuật toán âm học đo nhịp độ, tốc độ (WPM), "
                "khoảng lặng và đưa qua SenseVoice để trích xuất LID, SER, AED phục vụ đánh giá mức độ tự tin, lưu loát."
            )
            
        # Section 2.6 - Paragraph 516
        if "Thuật toán phân tích SenseVoice trong hệ thống: Audio service có ba tầng" in txt:
            p.text = (
                "Thuật toán phân tích âm thanh và SenseVoice trong hệ thống: Audio service được cấu trúc thành ba tầng phối hợp chặt chẽ: "
                "(1) Tầng 1 - Xử lý tín hiệu âm học & Thuật toán nhịp độ: Xử lý tuần tự các đoạn audio đã lưu bằng Faster-Whisper để có timestamp từng từ chính xác, "
                "từ đó thuật toán tính toán tốc độ nói WPM, phân tích khoảng lặng (pause detection) và độ biến thiên âm lượng (RMS/dB stddev). "
                "(2) Tầng 2 - SenseVoice Tag Extraction: Nạp mô hình SenseVoiceSmall, gọi generate với language='auto' để bóc tách LID (ngôn ngữ), SER (cảm xúc) và AED (sự kiện: tiếng thở, cười, ho, filler words). "
                "(3) Tầng 3 - Tổng hợp Đánh giá Phong thái (Observation Only): Kết hợp kết quả phân tích âm học và SenseVoice tags để xây dựng báo cáo phong thái giao tiếp, "
                "đánh giá mức độ tự tin, lưu loát hỗ trợ coaching, đồng thời tuân thủ nghiêm ngặt nguyên tắc không suy diễn chẩn đoán tâm lý."
            )
            
        # Section 2.6 - Paragraph 610, 611, 612
        if "Thiết kế hợp nhất SenseVoice - STT - DeepSeek: Luồng xử lý bắt đầu từ audio theo từng câu trả lời" in txt:
            p.text = (
                "Thiết kế hợp nhất SenseVoice - STT - DeepSeek: Kiến trúc phân định rõ vai trò của từng mô hình theo 2 giai đoạn: "
                "Ở giai đoạn realtime, AssemblyAI streaming (với Faster-Whisper fallback) đảm bảo cung cấp transcript tức thời để DeepSeek hiểu câu trả lời và sinh câu hỏi thích ứng liền mạch. "
                "Ở giai đoạn hậu kỳ, Faster-Whisper xử lý tuần tự để đảm bảo độ chính xác tối đa của văn bản, kết hợp thuật toán tính nhịp độ, tốc độ nói, khoảng lặng, "
                "và SenseVoice đảm nhiệm bóc tách LID, SER, AED nhằm đánh giá phong thái tự tin, lưu loát cho báo cáo coaching."
            )

    # 2. Update Table 1.5 (Table 7)
    t7 = doc.tables[7]
    t7.rows[1].cells[1].text = "ASR streaming thời gian thực"
    t7.rows[1].cells[2].text = "Transcript realtime giúp AI provider nhận diện câu trả lời và sinh câu hỏi thích ứng liền mạch."
    t7.rows[1].cells[3].text = "Phụ thuộc kết nối mạng; tự động fallback khi lỗi."
    
    t7.rows[2].cells[1].text = "ASR fallback realtime & Xử lý tuần tự hậu kỳ"
    t7.rows[2].cells[2].text = "1) Fallback transcript cục bộ khi streaming lỗi; 2) Xử lý tuần tự audio sau phỏng vấn để STT chính xác cao + thuật toán đo nhịp độ, tốc độ WPM, khoảng lặng."
    t7.rows[2].cells[3].text = "Tài nguyên tính toán server (chạy nền cục bộ)."
    
    t7.rows[3].cells[1].text = "LID + SER + AED (Observation only)"
    t7.rows[3].cells[2].text = "Nhận dạng ngôn ngữ, cảm xúc âm học, sự kiện âm thanh (tiếng thở, cười, ho, filler words); đánh giá mức độ tự tin, lưu loát cho coaching."
    t7.rows[3].cells[3].text = "Chỉ dùng quan sát mô tả; không chẩn đoán tâm lý hay tuyển/loại."

    doc.save(docx_path)
    print(f"Successfully updated {docx_path}")

def update_markdown_and_json():
    # 1. Update current_slides_analysis.md
    md_path = r"d:\project\InterV\Review\current_slides_analysis.md"
    if os.path.exists(md_path):
        with open(md_path, "r", encoding="utf-8") as f:
            content = f.read()
            
        old_slide15 = """## Slide 15

- Pipeline tiếng nói nhiều tầng
- 03
- AssemblyAI streaming • Faster-Whisper fallback • SenseVoice LID/SER/AED
- LUỒNG PHIÊN ÂM
- FALLBACK KHI STREAMING LỖI
- LUỒNG QUAN SÁT
- Microphone
- Audio đầu vào
- AssemblyAI
- Streaming transcript
- Transcript
- Theo thời gian thực
- SenseVoice
- LID / SER / AED
- Observation only
- Faster-Whisper
- Fallback cục bộ
- Transcript fallback
- Tiếp tục phiên"""

        new_slide15 = """## Slide 15

- Pipeline tiếng nói nhiều tầng: Realtime & Hậu kỳ phỏng vấn
- 03
- AssemblyAI streaming • Faster-Whisper fallback & batch • SenseVoice LID/SER/AED
- 1. LUỒNG REALTIME TRONG PHỎNG VẤN (TƯƠNG TÁC LIỀN MẠCH)
- 2. LUỒNG XỬ LÝ HẬU KỲ SAU PHỎNG VẤN (ĐÁNH GIÁ CHUYÊN SÂU)
- Microphone
- Thu âm AudioWorklet (16kHz PCM)
- AssemblyAI Streaming
- STT thời gian thực (WebSocket) [Lỗi → Faster-Whisper fallback]
- AI Provider (DeepSeek)
- Có text tức thời → Sinh câu hỏi thích ứng & TTS mượt mà, không rời rạc
- Faster-Whisper tuần tự
- STT chính xác cao + Thuật toán: nhịp độ, tốc độ (WPM), khoảng lặng
- SenseVoice (Observation only)
- LID (Ngôn ngữ) • SER (Cảm xúc) • AED (Tiếng thở, cười, filler words)
- Đánh giá phong thái
- Mức độ tự tin, lưu loát, nhịp điệu diễn đạt hỗ trợ coaching"""

        if old_slide15 in content:
            content = content.replace(old_slide15, new_slide15)
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(content)
            print(f"Successfully updated {md_path}")
            
    # 2. Update slide_notes_dump.json
    json_path = r"d:\project\InterV\Review\slide_notes_dump.json"
    if os.path.exists(json_path):
        with open(json_path, "r", encoding="utf-8") as f:
            notes = json.load(f)
            
        new_note_15 = (
            "[08:00 - 09:30] SLIDE 15: PIPELINE TIẾNG NÓI ĐA TẦNG (REALTIME & HẬU KỲ)\n\n"
            "Lời nói:\n"
            "\"Kính thưa Thầy Cô, pipeline tiếng nói của InterV được thiết kế chuyên biệt theo 2 pha rõ rệt:\n\n"
            "1. PHA 1: LUỒNG REALTIME TRONG PHỎNG VẤN (TƯƠNG TÁC MƯỢT MÀ, KHÔNG RỜI RẠC)\n"
            "• Luồng âm thanh từ microphone được AudioWorklet thu và truyền trực tiếp qua WebSocket tới AssemblyAI để thực hiện STT theo thời gian thực.\n"
            "• Cơ chế Fallback cục bộ: Khi gặp sự cố mạng hoặc lỗi API, hệ thống tự động fallback về model Faster-Whisper chạy trực tiếp trên Python backend.\n"
            "• Mục đích: Luôn bảo đảm có dữ liệu transcript tức thời cho AI Provider (DeepSeek) tiếp tục hiểu ngữ cảnh, nói tiếp qua TTS và sinh các câu hỏi thích ứng (Lookahead Engine) liền mạch, không làm câu hỏi bị rời rạc hay ngắt quãng phiên phỏng vấn.\n\n"
            "2. PHA 2: LUỒNG XỬ LÝ HẬU KỲ SAU PHỎNG VẤN (ĐÁNH GIÁ CHUYÊN SÂU & COACHING)\n"
            "• Sau khi phỏng vấn xong, toàn bộ các đoạn audio gốc đã lưu sẽ được xử lý tuần tự (sequential batch) bằng Faster-Whisper để tái tạo STT chính xác cao nhất.\n"
            "• Hệ thống áp dụng các thuật toán xử lý tín hiệu âm thanh để đo lường khách quan nhịp độ, tốc độ nói (WPM), và các khoảng trống/khoảng lặng (pause duration) trong câu nói.\n"
            "• Tiếp theo, audio được đưa qua mô hình SenseVoice theo nguyên tắc 'Observation only' để bóc tách: LID (ngôn ngữ), SER (cảm xúc giọng nói), và AED (sự kiện âm thanh: tiếng cười, tiếng ho, thở dài, từ đệm/filler words).\n"
            "• Dữ liệu này được tổng hợp để đánh giá khách quan mức độ tự tin, tính lưu loát, và phong thái diễn đạt nhằm hỗ trợ coaching cho ứng viên mà không suy diễn tâm lý tùy tiện.\""
        )
        if isinstance(notes, list):
            for item in notes:
                if item.get("slide_num") == 15:
                    item["notes"] = new_note_15
                    break
        elif isinstance(notes, dict):
            notes["15"] = new_note_15
            
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(notes, f, ensure_ascii=False, indent=2)
        print(f"Successfully updated {json_path}")

if __name__ == "__main__":
    update_pptx(r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v19.pptx")
    update_pptx(r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v18.pptx")
    update_kich_ban_docx()
    update_thesis_v14()
    update_markdown_and_json()
    print("ALL PPT AND WORD DOCUMENTS UPDATED SUCCESSFULLY!")
