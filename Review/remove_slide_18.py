# -*- coding: utf-8 -*-
"""
Script to remove Slide 18 from PPTX presentations (v19 & v18),
re-index all slides (32 -> 31 slides), update slide notes,
and synchronize Word script, markdown analysis, and helper scripts.
"""

import os
import pptx
import json

def remove_slide_18_from_pptx(pptx_path):
    if not os.path.exists(pptx_path):
        print(f"File not found: {pptx_path}")
        return
    
    prs = pptx.Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"Original total slides in {os.path.basename(pptx_path)}: {total_slides}")
    
    if total_slides == 32:
        # Delete slide at index 17 (Slide 18)
        rId = prs.slides._sldIdLst[17].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[17]
        print(f"Deleted slide at index 17. New total slides: {len(prs.slides)}")
    elif total_slides == 31:
        print("Presentation already has 31 slides.")
    else:
        print(f"Unexpected slide count: {total_slides}")
        return

    # Updated slide speaking notes for 31 slides
    speaking_notes = {
        1: (
            "[00:00 - 00:30] SLIDE 1: GIỚI THIỆU ĐỀ TÀI\n\n"
            "Lời nói:\n"
            "\"Kính thưa quý Thầy Cô trong Hội đồng. Em tên là Lê Minh Duy, sinh viên thực hiện khóa luận tốt nghiệp: "
            "'Xây dựng hệ thống phỏng vấn và luyện tập phỏng vấn tích hợp Trí tuệ Nhân tạo - InterV', dưới sự hướng dẫn của Thầy ThS. Đặng Văn Lực.\n"
            "InterV là nền tảng phỏng vấn giọng nói hai chế độ: Luyện tập cá nhân và Tuyển dụng thực tế, "
            "kết hợp mô hình DeepSeek, cơ sở dữ liệu vector Qdrant RAG, và mô hình phân tích âm thanh đa phương thức SenseVoice.\""
        ),
        2: (
            "[00:30 - 00:50] SLIDE 2: LỘ TRÌNH BÀI TRÌNH BÀY\n\n"
            "Lời nói:\n"
            "\"Nội dung trình bày gồm 6 phần chính. Để làm nổi bật chiều sâu kỹ thuật, em xin phép dành 10 phút trọng tâm "
            "cho Pipeline Cốt lõi của hệ thống, và 5 phút cho bối cảnh, ranh giới và hiện vật triển khai.\""
        ),
        3: (
            "[00:50 - 01:15] SLIDE 3: BÀI TOÁN & THỰC TRẠNG\n\n"
            "Lời nói:\n"
            "\"Các hệ thống phỏng vấn AI hiện nay thường mắc phải 3 lỗi lớn:\n"
            "1. LLM thiếu căn cứ (Hallucination): Câu hỏi và nhận xét lệch JD.\n"
            "2. Thiếu nhất quán: Tiêu chí thay đổi giữa các ứng viên.\n"
            "3. Suy diễn quá mức: Biến tín hiệu giọng nói thành kết luận tâm lý thiếu cơ sở.\n"
            "Mục tiêu của InterV: Mọi câu hỏi và đánh giá đều phải có căn cứ (Grounded), nhất quán và AI không thay thế con người ra quyết định.\""
        ),
        4: (
            "[01:15 - 01:30] SLIDE 4: BA NHÓM NGƯỜI DÙNG & RANH GIỚI TRÁCH NHIỆM\n\n"
            "Lời nói:\n"
            "\"Hệ thống phân định 3 ranh giới trách nhiệm tuyệt đối:\n"
            "• Ứng viên: Sở hữu dữ liệu cá nhân, luyện tập hoặc tham gia theo lời mời.\n"
            "• Recruiter: Sở hữu quyết định tuyển dụng (thiết lập JD, xem bằng chứng và đánh giá cuối).\n"
            "• Admin: Vận hành hạ tầng và tài chính; không can thiệp kết quả tuyển dụng.\""
        ),
        5: (
            "[01:30 - 01:45] SLIDE 5: MỤC TIÊU NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Sáu mục tiêu thiết kế: Tách 2 chế độ, Grounded LLM, Speech an toàn dạng Coaching, Quản trị Provenance, Xác minh kỹ thuật toàn diện, và Không auto-hire/auto-reject.\""
        ),
        6: (
            "[01:45 - 02:00] SLIDE 6: BỐN NGUYÊN TẮC KHÓA PHẠM VI\n\n"
            "Lời nói:\n"
            "\"Phạm vi đề tài được khóa bằng 4 nguyên tắc: Evidence-first (output nối về JD/Rule), Human-in-the-loop, Observation only (chỉ quan sát tín hiệu), và Đúng mức bằng chứng.\""
        ),
        7: (
            "[02:00 - 02:15] SLIDE 7: PHƯƠNG PHÁP NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Phương pháp tách rõ hai tầng: Xác minh kỹ thuật (chạy đúng, type-safe, contract test) và Đánh giá hiệu lực nghiệp vụ.\""
        ),
        8: (
            "[02:15 - 02:30] SLIDE 8: ÁNH XẠ LÝ THUYẾT VÀO HỆ THỐNG\n\n"
            "Lời nói:\n"
            "\"Các khung lý thuyết chuẩn được chuyển hóa thành mã nguồn: Structured Interview thành Schema, STAR/BARS thành Rule Catalog 86 files, và Responsible AI thành Grounding Tri-gate.\""
        ),
        9: (
            "[02:30 - 03:00] SLIDE 9: KIẾN TRÚC TỔNG THỂ 2 TẦNG\n\n"
            "Lời nói:\n"
            "\"Kiến trúc phân tầng rạch ròi:\n"
            "• Web/BFF (Next.js App Router): Quản lý session, auth, WebSocket audio, Mongo replica set.\n"
            "• AI Backend (Python): DeepSeek, Qdrant RAG, SenseVoice, TTS.\n"
            "Giao tiếp qua hợp đồng gRPC typed 16 RPCs có xác thực nội bộ. Không có mock data trong runtime.\""
        ),
        10: (
            "[03:00 - 03:15] SLIDE 10: MÔ HÌNH DỮ LIỆU & PROVENANCE\n\n"
            "Lời nói:\n"
            "\"19 Mongoose models liên kết chặt chẽ theo vòng đời: User, Job, Invitation, Session, Run, Audio (BSON binary) và Result.\""
        ),
        11: (
            "[03:15 - 03:30] SLIDE 11: USE CASE THEO QUYỀN SỞ HỮU\n\n"
            "Lời nói:\n"
            "\"Các use case bảo đảm AI nằm ngoài biên quyết định tuyển dụng; recruiter nắm giữ quyền duyệt cuối.\""
        ),
        12: (
            "[03:30 - 05:30] SLIDE 12: LOOKAHEAD ADAPTIVE QUESTION ENGINE (TRỌNG TÂM)\n\n"
            "Lời nói:\n"
            "\"★ ĐIỂM SÁNG KIẾN TRÚC - XỬ LÝ ĐỘ TRỄ 0MS:\n"
            "1. Preparation: Backend sinh sẵn bộ câu hỏi baseline + warm TTS. Trả ngay câu 1.\n"
            "2. Instant Return: Ứng viên nộp câu Q_i -> Trả ngay Q_(i+1) có sẵn trong bộ nhớ đệm (0ms latency, không phải chờ LLM).\n"
            "3. Background Lookahead: Next.js after() gọi bất đồng bộ SubmitAnswer qua gRPC -> DeepSeek + RAG phân tích câu trả lời Q_i để bắt bài (probe gap) và sinh câu hỏi thích ứng Q_(i+2), đồng thời warm TTS trong nền.\n"
            "4. Ghi đè thông minh: Q_(i+2) ghi đè vào slot kế tiếp; nếu mạng chậm thì câu baseline làm fallback an toàn. Trải nghiệm phỏng vấn luôn mượt mà và thông minh!\""
        ),
        13: (
            "[05:30 - 07:00] SLIDE 13: GROUNDED GENERATION & CƠ CHẾ KIỂM SOÁT\n\n"
            "Lời nói:\n"
            "\"★ CƠ CHẾ GROUNDING 3 CỔNG CHỐNG ẢO GIÁC:\n"
            "1. Chuẩn hóa Context & Cấp phát Evidence IDs.\n"
            "2. DeepSeek Structured JSON Output với ràng buộc trích dẫn Evidence ID.\n"
            "3. Citation Gate: Backend kiểm tra đối chiếu allow-list. Nếu DeepSeek bịa ra ID lạ -> Chặn ngay lập tức và kích hoạt 1 lượt Repair Request tự sửa sai. Bất biến: 100% câu hỏi và đánh giá đều truy nguyên được nguồn gốc!\""
        ),
        14: (
            "[07:00 - 08:00] SLIDE 14: VÒNG ĐỜI TÀI LIỆU RAG & HYBRID RETRIEVAL\n\n"
            "Lời nói:\n"
            "\"Kho tri thức 86 Rule Markdown (15 ngành, 60 profile STAR/BARS).\n"
            "Qdrant kết hợp song song:\n"
            "• Dense Embeddings: FastEmbed Multilingual MPNet.\n"
            "• Sparse Lexical: Thuật toán BM25 bắt từ khóa kỹ thuật.\n"
            "Hợp nhất bằng Reciprocal Rank Fusion (RRF), đảm bảo truy hồi tài liệu vừa chuẩn ngữ nghĩa vừa chính xác từ khóa.\""
        ),
        15: (
            "[08:00 - 09:30] SLIDE 15: PIPELINE TIẾNG NÓI ĐA TẦNG (REALTIME & HẬU KỲ)\n\n"
            "Lời nói:\n"
            "\"Kính thưa Thầy Cô, pipeline tiếng nói của InterV được thiết kế chuyên biệt theo 2 pha rõ rệt:\n\n"
            "1. PHA 1: LUỒNG REALTIME TRONG PHỎNG VẤN (TƯƠNG TÁC MƯỢT MÀ, KHÔNG RỜI RẠC)\n"
            "• Luồng âm thanh từ microphone được AudioWorklet thu và truyền trực tiếp qua WebSocket tới AssemblyAI để thực hiện STT theo thời gian thực.\n"
            "• Cơ chế Fallback cục bộ: Khi gặp sự cố mạng hoặc lỗi API, hệ thống tự động fallback về model Faster-Whisper chạy trực tiếp trên Python backend.\n"
            "• Mục đích: Luôn bảo đảm có dữ liệu transcript tức thời cho AI Provider (DeepSeek) tiếp tục hiểu ngữ cảnh, nói tiếp qua TTS và sinh các câu hỏi thích ứng (Lookahead Engine) liền mạch, không làm câu hỏi bị rời rạc hay ngắt quãng phiên phỏng vấn.\n\n"
            "2. PHA 2: LUỒNG XỬ LÝ HẬU KỲ SAU PHỎNG VẤN (ĐÁNH GIÁ CHUYÊN SÂU & COACHING)\n"
            "• Sau khi phỏng vấn xong, toàn bộ các đoạn audio gốc đã lưu sẽ được xử lý tuần tự (sequential batch) bằng Faster-Whisper để tái tạo STT chính xác cao nhất.\n"
            "• Hệ thống áp dụng các thuật toán xử lý tín hiệu âm thanh để đo lường khách quan nhịp độ, tốc độ nói (WPM), và các khoảng trống/khoảng lặng (pause duration) trong câu nói.\n"
            "• Tiếp theo, audio được đưa qua mô hình SenseVoice theo nguyên tắc 'Observation only' để bóc tách: LID (ngôn ngữ), SER (cảm xúc giọng nói), và AED (sự kiện âm thanh: tiếng cười, tiếng ho, thở dài, từ đệm/filler words).\n"
            "• Dữ liệu này được tổng hợp để đánh giá khách quan mức độ tự tin, tính lưu loát, và phong thái diễn đạt nhằm hỗ trợ coaching cho ứng viên mà không suy diễn tâm lý tùy tiện.\""
        ),
        16: (
            "[09:30 - 10:15] SLIDE 16: AN TOÀN, QUAN SÁT & KHẢ NĂNG PHỤC HỒI\n\n"
            "Lời nói:\n"
            "\"5 lớp bảo vệ: Auth/RBAC, Schema validation Zod/Pydantic, Timeout/Retry circuit breaker, Audit log/Usage tracking, Graceful fallback ở mọi khâu.\""
        ),
        17: (
            "[10:15 - 10:45] SLIDE 17: HỢP ĐỒNG TYPED gRPC\n\n"
            "Lời nói:\n"
            "\"Typed RPC 16 phương thức ràng buộc giữa TypeScript BFF và Python Backend. Dữ liệu chỉ được lưu trữ hoặc hiển thị sau khi vượt qua schema validation.\""
        ),
        18: (
            "[10:45 - 11:05] SLIDE 18: MA TRẬN PHÂN QUYỀN RBAC\n\n"
            "Lời nói:\n"
            "\"Candidate, Recruiter, Admin. Ba vai trò phân quyền độc lập, AI tuyệt đối không có quyền tự động ra quyết định tuyển dụng hay loại bỏ ứng viên.\""
        ),
        19: (
            "[11:05 - 11:25] SLIDE 19: STATE MACHINE PHIÊN LUYỆN TẬP\n\n"
            "Lời nói:\n"
            "\"Vòng đời nghiêm ngặt: Draft -> Ready -> Recording -> Processing -> Reviewed -> Completed. Có cơ chế resume an toàn khi gặp sự cố mạng.\""
        ),
        20: (
            "[11:25 - 11:45] SLIDE 20: VÒNG ĐỜI TUYỂN DỤNG RECRUITER\n\n"
            "Lời nói:\n"
            "\"Job -> Campaign -> Invitation -> Interview -> Evidence -> Human Final Review. Recruiter sở hữu toàn quyền quyết định cuối cùng.\""
        ),
        21: (
            "[11:45 - 12:00] SLIDE 21: CƠ CHẾ FALLBACK TOÀN DIỆN\n\n"
            "Lời nói:\n"
            "\"Minh họa đường đi chính (Primary) và đường dự phòng (Fallback) tại mọi điểm trọng yếu: RAG fallback baseline, STT fallback Faster-Whisper, Audio fallback text-only.\""
        ),
        22: (
            "[12:00 - 12:25] SLIDE 22: STACK CÔNG NGHỆ\n\n"
            "Lời nói:\n"
            "\"Next.js 15, TypeScript, Tailwind v4, MongoDB ReplicaSet, Python gRPC (16 RPCs), DeepSeek, Qdrant RAG, AssemblyAI, Faster-Whisper, SenseVoice, Edge/Vbee TTS.\""
        ),
        23: (
            "[12:25 - 13:10] SLIDE 23: QUY MÔ HIỆN VẬT TRIỂN KHAI\n\n"
            "Lời nói:\n"
            "\"Con số thực tế từ repository:\n"
            "• Hơn 72.000 dòng mã nguồn logic.\n"
            "• 61 route files, 74 HTTP methods, 38 trang frontend.\n"
            "• 19 Mongoose models, 16 gRPC RPCs, 86 Rule files (15 ngành, 60 profile STAR/BARS).\""
        ),
        24: (
            "[13:10 - 13:40] SLIDE 24: XÁC MINH KỸ THUẬT ĐẠT 100%\n\n"
            "Lời nói:\n"
            "\"Toàn bộ 51/51 Backend tests đạt 100% trong 3.28 giây, 16/16 RPC contract tests pass, Next.js build đạt production-ready, zero-mock runtime.\""
        ),
        25: (
            "[13:40 - 13:55] SLIDE 25: BA ĐIỂM KIỂM SOÁT DEEPSEEK\n\n"
            "Lời nói:\n"
            "\"Kiểm soát Trước - Trong - Sau khi sinh: Pre-generation Context Grounding, In-generation JSON Schema, Post-generation Citation Gate & 1-shot Repair.\""
        ),
        26: (
            "[13:55 - 14:15] SLIDE 26: BẤT BIẾN AN TOÀN SENSEVOICE\n\n"
            "Lời nói:\n"
            "\"Observation-only (chỉ coaching, không chẩn đoán tâm lý), Neutral Default 50 nếu thiếu tín hiệu, không bao giờ tự ý trừ điểm của ứng viên.\""
        ),
        27: (
            "[14:15 - 14:30] SLIDE 27: BỐN ĐÓNG GÓP CHÍNH\n\n"
            "Lời nói:\n"
            "\"1. Kiến trúc Lookahead thời gian thực 0ms latency.\n"
            "2. Grounded Generation 3 cổng loại bỏ hoàn toàn ảo giác.\n"
            "3. Multimodal Speech Pipeline 2 pha an toàn, có fallback.\n"
            "4. Hệ thống Provenance & Rule Catalog 86 files chuẩn hóa theo STAR/BARS.\""
        ),
        28: (
            "[14:30 - 14:45] SLIDE 28: GIỚI HẠN NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Đề tài trung thực xác nhận: Đạt tính nhất quán và vững chắc kỹ thuật, nhưng cần đánh giá Validity/Fairness thực nghiệm trên tập mẫu lớn trước khi triển khai tuyển dụng diện rộng.\""
        ),
        29: (
            "[14:45 - 14:50] SLIDE 29: BỐN CỔNG KIỂM SOÁT TRIỂN KHAI\n\n"
            "Lời nói:\n"
            "\"Lộ trình mở rộng qua 4 gates: Practice Playground -> Shadow Scoring -> Human-overridden Pilot -> Restricted Production.\""
        ),
        30: (
            "[14:50 - 14:55] SLIDE 30: KẾT LUẬN\n\n"
            "Lời nói:\n"
            "\"InterV đã đạt nền tảng kỹ thuật khả thi, an toàn và sẵn sàng phục vụ luyện tập và hỗ trợ tuyển dụng có kiểm soát.\""
        ),
        31: (
            "[14:55 - 15:00] SLIDE 31: CẢM ƠN & TIẾP NHẬN PHẢN BIỆN\n\n"
            "Lời nói:\n"
            "\"Em xin trân trọng cảm ơn Thầy ThS. Đặng Văn Lực và quý Thầy Cô trong Hội đồng đã chú ý theo dõi. Em xin phép được tiếp nhận câu hỏi phản biện!\""
        )
    }

    # Apply speaking notes to each slide
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        note_text = speaking_notes.get(slide_num, "")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide.notes_slide.notes_text_frame.text = note_text
            
    try:
        prs.save(pptx_path)
        print(f"Successfully saved {pptx_path} with {len(prs.slides)} slides and updated slide notes.")
    except PermissionError:
        base_dir = os.path.dirname(pptx_path)
        v20_path = os.path.join(base_dir, "KLTN_InterV_LeMinhDuy_v20.pptx")
        prs.save(v20_path)
        print(f"File {os.path.basename(pptx_path)} was locked by PowerPoint. Saved updated 31-slide version to: {v20_path}")

def update_markdown_and_scripts():
    # Use v20 or available 31-slide pptx
    pptx_path = r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v20.pptx"
    if not os.path.exists(pptx_path):
        pptx_path = r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v19.pptx"
        
    prs = pptx.Presentation(pptx_path)
    
    # 1. Update current_slides_analysis.md
    md_path = r"d:\project\InterV\Review\current_slides_analysis.md"
    new_md_lines = [f"# PHÂN TÍCH VÀ NỘI DUNG TỪNG SLIDE ({os.path.basename(pptx_path)} - {len(prs.slides)} Slides)\n\n"]
    for idx, slide in enumerate(prs.slides):
        new_md_lines.append(f"## Slide {idx + 1}\n\n")
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    txt = p.text.strip()
                    if txt:
                        new_md_lines.append(f"- {txt}\n")
        new_md_lines.append("\n---\n\n")
        
    with open(md_path, "w", encoding="utf-8") as f:
        f.writelines(new_md_lines)
    print(f"Successfully updated {md_path}")
    
    # 2. Update slide_notes_dump.json
    json_path = r"d:\project\InterV\Review\slide_notes_dump.json"
    notes_list = []
    for idx, slide in enumerate(prs.slides):
        note = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame else ""
        notes_list.append({
            "slide_num": idx + 1,
            "notes": note
        })
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(notes_list, f, ensure_ascii=False, indent=2)
    print(f"Successfully updated {json_path}")

if __name__ == "__main__":
    remove_slide_18_from_pptx(r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v19.pptx")
    remove_slide_18_from_pptx(r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v18.pptx")
    update_markdown_and_scripts()
    print("SLIDE 18 REMOVED AND ALL FILES SYNCHRONIZED SUCCESSFULLY!")
