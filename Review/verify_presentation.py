import pptx

with open(r"d:\project\InterV\Review\verify_result.txt", "w", encoding="utf-8") as out:
    for path in [r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v18.pptx", r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v19.pptx"]:
        prs = pptx.Presentation(path)
        s24 = prs.slides[23]
        s25 = prs.slides[24]
        t24 = [p.text for sh in s24.shapes if sh.has_text_frame for p in sh.text_frame.paragraphs if p.text.strip()]
        t25 = [p.text for sh in s25.shapes if sh.has_text_frame for p in sh.text_frame.paragraphs if p.text.strip()]
        notes12 = prs.slides[11].notes_slide.notes_text_frame.text if prs.slides[11].has_notes_slide else ""
        out.write(f"File: {path}\n")
        out.write(f"  Total slides: {len(prs.slides)}\n")
        out.write(f"  Slide 24 text snippet: {t24}\n")
        out.write(f"  Slide 25 text snippet: {t25}\n")
        out.write(f"  Slide 12 Notes snippet: {notes12[:150]}\n")
        out.write("-" * 50 + "\n")

    out.write("Verification complete!\n")
print("Done writing verify_result.txt")
