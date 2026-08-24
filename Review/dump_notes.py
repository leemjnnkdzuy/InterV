import pptx
import json

prs = pptx.Presentation(r'd:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v18.pptx')
notes_data = []

for idx, slide in enumerate(prs.slides):
    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    notes_data.append({
        "slide_num": idx + 1,
        "notes": notes
    })

with open(r'd:\project\InterV\Review\slide_notes_dump.json', 'w', encoding='utf-8') as f:
    json.dump(notes_data, f, ensure_ascii=False, indent=2)

print("Dumped slide_notes_dump.json successfully")
