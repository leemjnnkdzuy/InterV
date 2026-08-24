# -*- coding: utf-8 -*-
"""
Master Build Script for InterV 37-Slide Presentation:
1. Removes old Slide 18 ('Phân rã module theo ranh giới nghiệp vụ').
2. Inserts '#xx' (#01 .. #06) top-left badge on all content slides.
3. Creates 6 Section Divider slides with safe unique part names.
4. Moves dividers into exact section transition positions.
5. Injects speaking notes for all 37 slides.
6. Saves to KLTN_InterV_LeMinhDuy_v20.pptx & v18.pptx.
7. Updates current_slides_analysis.md, slide_notes_dump.json, and Word script.
"""

import os
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.opc.packuri import PackURI
from pptx.parts.slide import SlidePart
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
import json

def add_slide_safe(prs, slide_layout, num):
    """Add a slide with a guaranteed unique partname to prevent collision."""
    partname = PackURI(f"/ppt/slides/slide_divider_{num}.xml")
    slide_part = SlidePart.new(partname, prs.part.package, slide_layout.part)
    rId = prs.part.relate_to(slide_part, RT.SLIDE)
    slide = slide_part.slide
    slide.shapes.clone_layout_placeholders(slide_layout)
    prs.slides._sldIdLst.add_sldId(rId)
    return slide

def build_master_deck():
    src_path = r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v19.pptx"
    prs = pptx.Presentation(src_path)
    print(f"Loaded {os.path.basename(src_path)} with {len(prs.slides)} slides.")
    
    # 0. Set official thesis title on Slide 1
    slide1 = prs.slides[0]
    for s in slide1.shapes:
        if s.has_text_frame and s.top < 2000000 and s.height > 1000000:
            tf = s.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "NGHIÊN CỨU THIẾT KẾ VÀ PHÁT TRIỂN KIẾN TRÚC HỆ THỐNG PHỎNG VẤN NĂNG LỰC TỰ ĐỘNG TÍCH HỢP XỬ LÝ GIỌNG NÓI ĐA TẦNG, TRUY XUẤT TRI THỨC TĂNG CƯỜNG (RAG) VÀ MÔ HÌNH NGÔN NGỮ LỚN"
            p.font.name = "Cambria"
            p.font.size = Pt(23.25)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
            p.alignment = PP_ALIGN.CENTER

    # 1. Delete old slide 18 (index 17)
    rId = prs.slides._sldIdLst[17].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[17]
    print(f"Pristine content slides: {len(prs.slides)}") # 31 slides
    
    sections_def = [
        {
            "sec_num": "01",
            "tag": "#01",
            "title": "Bài toán",
            "subtitle": "Bối cảnh, đối tượng, mục tiêu và 4 nguyên tắc khóa phạm vi",
            "orig_idx": 2,
            "slide_indices": [2, 3, 4, 5]
        },
        {
            "sec_num": "02",
            "tag": "#02",
            "title": "Phương pháp",
            "subtitle": "Nền tảng lý thuyết, xác minh kỹ thuật và ánh xạ vào hệ thống",
            "orig_idx": 6,
            "slide_indices": [6, 7]
        },
        {
            "sec_num": "03",
            "tag": "#03",
            "title": "Thiết kế",
            "subtitle": "Kiến trúc 2 tầng, Lookahead 0ms, RAG Hybrid Grounding, Pipeline tiếng nói 2 pha & Typed gRPC",
            "orig_idx": 8,
            "slide_indices": [8, 9, 10, 11, 12, 13, 14, 15, 16]
        },
        {
            "sec_num": "04",
            "tag": "#04",
            "title": "Sản phẩm",
            "subtitle": "Ma trận RBAC, State Machine luyện tập, Vòng đời tuyển dụng Recruiter & Graceful Fallback",
            "orig_idx": 17,
            "slide_indices": [17, 18, 19, 20, 21]
        },
        {
            "sec_num": "05",
            "tag": "#05",
            "title": "Bằng chứng",
            "subtitle": "Quy mô hiện vật 72.000+ LOC, 51/51 Backend tests đạt 100%, 3 điểm kiểm soát LLM & Bất biến SenseVoice",
            "orig_idx": 22,
            "slide_indices": [22, 23, 24, 25]
        },
        {
            "sec_num": "06",
            "tag": "#06",
            "title": "Kết luận",
            "subtitle": "4 đóng góp then chốt, giới hạn nghiên cứu, 4 cổng kiểm soát triển khai & Lời cảm ơn",
            "orig_idx": 26,
            "slide_indices": [26, 27, 28, 29, 30]
        }
    ]
    
    # 2. Add top-left #xx badge to every content slide (aligned to left, vertically centered in top strip)
    for sec in sections_def:
        tag_text = sec["tag"]
        for s_idx in sec["slide_indices"]:
            slide = prs.slides[s_idx]
            
            # Remove old placeholder numbers (like "01", "02", "03" in top-right)
            for shape in slide.shapes:
                if shape.has_text_frame and shape.top < 1200000 and shape.left > 6000000:
                    t = shape.text_frame.text.strip()
                    if t in ["01", "02", "03", "04", "05", "06", "#01", "#02", "#03", "#04", "#05", "#06"]:
                        shape.text_frame.clear()
            
            # Check for existing badge
            existing_badge = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape.top < 450000 and shape.left < 2500000:
                    if shape.text_frame.text.strip().startswith("#"):
                        existing_badge = shape
                        break
                        
            if existing_badge:
                existing_badge.left = 274100
                existing_badge.top = 0
                existing_badge.width = 1000000
                existing_badge.height = 274100
                tf = existing_badge.text_frame
                tf.word_wrap = False
                tf.margin_left = 0
                tf.margin_top = 0
                tf.margin_right = 0
                tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.clear()
                p = tf.paragraphs[0]
                p.text = tag_text
                p.font.name = "Montserrat"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
                p.alignment = PP_ALIGN.LEFT
            else:
                tx_box = slide.shapes.add_textbox(left=274100, top=0, width=1000000, height=274100)
                tf = tx_box.text_frame
                tf.word_wrap = False
                tf.margin_left = 0
                tf.margin_top = 0
                tf.margin_right = 0
                tf.margin_bottom = 0
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = tag_text
                p.font.name = "Montserrat"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
                p.alignment = PP_ALIGN.LEFT

    # 3. Create and insert 6 Section Divider slides backwards
    for div_idx, sec in enumerate(reversed(sections_def)):
        div_num = 6 - div_idx
        d_slide = add_slide_safe(prs, prs.slide_layouts[1], div_num) # SECTION_HEADER
        
        for shape in d_slide.shapes:
            if shape.is_placeholder:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 2 or shape.top < 2000000:
                    shape.left = 1072000
                    shape.width = 7000000
                    shape.top = 1300000
                    shape.height = 800000
                    tf = shape.text_frame
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = sec["sec_num"]
                    p.font.name = "Cambria"
                    p.font.size = Pt(44)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0x4B, 0x8D, 0x17)
                    p.alignment = PP_ALIGN.CENTER
                elif ph_idx == 0 or (shape.top >= 2000000 and shape.top < 2900000):
                    shape.left = 1072000
                    shape.width = 7000000
                    shape.top = 2100000
                    shape.height = 800000
                    tf = shape.text_frame
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = sec["title"].upper()
                    p.font.name = "Cambria"
                    p.font.size = Pt(28)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0x11, 0x14, 0x10)
                    p.alignment = PP_ALIGN.CENTER
                elif ph_idx == 1 or shape.top >= 2900000:
                    shape.left = 572000
                    shape.width = 8000000
                    shape.top = 3050000
                    shape.height = 900000
                    tf = shape.text_frame
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0
                    tf.word_wrap = True
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = sec["subtitle"]
                    p.font.name = "Montserrat"
                    p.font.size = Pt(13)
                    p.font.color.rgb = RGBColor(0x52, 0x60, 0x52)
                    p.alignment = PP_ALIGN.CENTER
                    
        elem = prs.slides._sldIdLst[-1]
        prs.slides._sldIdLst.insert(sec["orig_idx"], elem)

    print(f"Total slides after section divider insertion: {len(prs.slides)}") # Exactly 37 slides

    # 4. Populate speaking notes for all 37 slides
    speaking_notes_37 = {
        1: (
            "[00:00 - 00:30] SLIDE 1: GIỚI THIỆU ĐỀ TÀI\n\n"
            "Lời nói:\n"
            "\"Kính thưa quý Thầy Cô trong Hội đồng chấm khóa luận tốt nghiệp. "
            "Em tên là Lê Minh Duy, sinh viên thực hiện đề tài khóa luận: "
            "'Nghiên cứu thiết kế và phát triển kiến trúc hệ thống phỏng vấn năng lực tự động tích hợp xử lý giọng nói đa tầng, truy xuất tri thức tăng cường (RAG) và mô hình ngôn ngữ lớn' (InterV), "
            "dưới sự hướng dẫn tận tình của Thầy ThS. Đặng Văn Lực.\n"
            "InterV là nền tảng phỏng vấn giọng nói hai chế độ: Luyện tập cá nhân và Tuyển dụng thực tế, "
            "kết hợp mô hình DeepSeek, cơ sở dữ liệu vector Qdrant RAG, và mô hình phân tích âm thanh đa phương thức SenseVoice.\""
        ),
        2: (
            "[00:30 - 00:50] SLIDE 2: LỘ TRÌNH BÀI TRÌNH BÀY\n\n"
            "Lời nói:\n"
            "\"Nội dung trình bày gồm 6 phần chính: 01 Bài toán, 02 Phương pháp, 03 Thiết kế, 04 Sản phẩm, 05 Bằng chứng, và 06 Kết luận. "
            "Để làm nổi bật chiều sâu kỹ thuật, em xin phép dành 10 phút trọng tâm cho Pipeline Cốt lõi của hệ thống, và 5 phút cho bối cảnh, ranh giới và hiện vật triển khai.\""
        ),
        3: (
            "[00:50 - 00:55] PHẦN 01: BÀI TOÁN & MỤC TIÊU\n\n"
            "Lời nói:\n"
            "\"Sau đây, em xin phép bắt đầu với Phần 1: Bài toán thực tế, ba nhóm người dùng và các ranh giới thiết kế cốt lõi của InterV.\""
        ),
        4: (
            "[00:55 - 01:15] SLIDE 4: BÀI TOÁN & THỰC TRẠNG\n\n"
            "Lời nói:\n"
            "\"Các hệ thống phỏng vấn AI hiện nay thường mắc phải 3 lỗi lớn:\n"
            "1. LLM thiếu căn cứ (Hallucination): Câu hỏi và nhận xét lệch JD.\n"
            "2. Thiếu nhất quán: Tiêu chí thay đổi giữa các ứng viên.\n"
            "3. Suy diễn quá mức: Biến tín hiệu giọng nói thành kết luận tâm lý thiếu cơ sở.\n"
            "Mục tiêu của InterV: Mọi câu hỏi và đánh giá đều phải có căn cứ (Grounded), nhất quán và AI không thay thế con người ra quyết định.\""
        ),
        5: (
            "[01:15 - 01:30] SLIDE 5: BA NHÓM NGƯỜI DÙNG & RANH GIỚI TRÁCH NHIỆM\n\n"
            "Lời nói:\n"
            "\"Hệ thống phân định 3 ranh giới trách nhiệm tuyệt đối:\n"
            "• Ứng viên: Sở hữu dữ liệu cá nhân, luyện tập hoặc tham gia theo lời mời.\n"
            "• Recruiter: Sở hữu quyết định tuyển dụng (thiết lập JD, xem bằng chứng và đánh giá cuối).\n"
            "• Admin: Vận hành hạ tầng và tài chính; không can thiệp kết quả tuyển dụng.\""
        ),
        6: (
            "[01:30 - 01:45] SLIDE 6: MỤC TIÊU NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Sáu mục tiêu thiết kế: Tách 2 chế độ, Grounded LLM, Speech an toàn dạng Coaching, Quản trị Provenance, Xác minh kỹ thuật toàn diện, và Không auto-hire/auto-reject.\""
        ),
        7: (
            "[01:45 - 02:00] SLIDE 7: BỐN NGUYÊN TẮC KHÓA PHẠM VI\n\n"
            "Lời nói:\n"
            "\"Phạm vi đề tài được khóa bằng 4 nguyên tắc: Evidence-first (output nối về JD/Rule), Human-in-the-loop, Observation only (chỉ quan sát tín hiệu), và Đúng mức bằng chứng.\""
        ),
        8: (
            "[02:00 - 02:05] PHẦN 02: PHƯƠNG PHÁP NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Tiếp theo là Phần 2: Phương pháp nghiên cứu và cách thức chuyển hóa các khung lý thuyết chuẩn thành mã nguồn.\""
        ),
        9: (
            "[02:05 - 02:20] SLIDE 9: PHƯƠNG PHÁP NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Phương pháp tách rõ hai tầng: Xác minh kỹ thuật (chạy đúng, type-safe, contract test) và Đánh giá hiệu lực nghiệp vụ.\""
        ),
        10: (
            "[02:20 - 02:35] SLIDE 10: ÁNH XẠ LÝ THUYẾT VÀO HỆ THỐNG\n\n"
            "Lời nói:\n"
            "\"Các khung lý thuyết chuẩn được chuyển hóa thành mã nguồn: Structured Interview thành Schema, STAR/BARS thành Rule Catalog 86 files, và Responsible AI thành Grounding Tri-gate.\""
        ),
        11: (
            "[02:35 - 02:40] PHẦN 03: THIẾT KẾ KIẾN TRÚC & PIPELINE\n\n"
            "Lời nói:\n"
            "\"Sau đây là Phần 3 - Trọng tâm bài thuyết trình: Thiết kế Kiến trúc và Pipeline Cốt lõi của InterV.\""
        ),
        12: (
            "[02:40 - 03:05] SLIDE 12: KIẾN TRÚC TỔNG THỂ 2 TẦNG\n\n"
            "Lời nói:\n"
            "\"Kiến trúc phân tầng rạch ròi:\n"
            "• Web/BFF (Next.js App Router): Quản lý session, auth, WebSocket audio, Mongo replica set.\n"
            "• AI Backend (Python): DeepSeek, Qdrant RAG, SenseVoice, TTS.\n"
            "Giao tiếp qua hợp đồng gRPC typed 16 RPCs có xác thực nội bộ. Không có mock data trong runtime.\""
        ),
        13: (
            "[03:05 - 03:20] SLIDE 13: MÔ HÌNH DỮ LIỆU & PROVENANCE\n\n"
            "Lời nói:\n"
            "\"19 Mongoose models liên kết chặt chẽ theo vòng đời: User, Job, Invitation, Session, Run, Audio (BSON binary) và Result.\""
        ),
        14: (
            "[03:20 - 03:35] SLIDE 14: USE CASE THEO QUYỀN SỞ HỮU\n\n"
            "Lời nói:\n"
            "\"Các use case bảo đảm AI nằm ngoài biên quyết định tuyển dụng; recruiter nắm giữ quyền duyệt cuối.\""
        ),
        15: (
            "[03:35 - 05:30] SLIDE 15: LOOKAHEAD ADAPTIVE QUESTION ENGINE (TRỌNG TÂM)\n\n"
            "Lời nói:\n"
            "\"★ ĐIỂM SÁNG KIẾN TRÚC - XỬ LÝ ĐỘ TRỄ 0MS:\n"
            "1. Preparation: Backend sinh sẵn bộ câu hỏi baseline + warm TTS. Trả ngay câu 1.\n"
            "2. Instant Return: Ứng viên nộp câu Q_i -> Trả ngay Q_(i+1) có sẵn trong bộ nhớ đệm (0ms latency, không phải chờ LLM).\n"
            "3. Background Lookahead: Next.js after() gọi bất đồng bộ SubmitAnswer qua gRPC -> DeepSeek + RAG phân tích câu trả lời Q_i để bắt bài (probe gap) và sinh câu hỏi thích ứng Q_(i+2), đồng thời warm TTS trong nền.\n"
            "4. Ghi đè thông minh: Q_(i+2) ghi đè vào slot kế tiếp; nếu mạng chậm thì câu baseline làm fallback an toàn. Trải nghiệm phỏng vấn luôn mượt mà và thông minh!\""
        ),
        16: (
            "[05:30 - 07:00] SLIDE 16: GROUNDED GENERATION & CƠ CHẾ KIỂM SOÁT\n\n"
            "Lời nói:\n"
            "\"★ CƠ CHẾ GROUNDING 3 CỔNG CHỐNG ẢO GIÁC:\n"
            "1. Chuẩn hóa Context & Cấp phát Evidence IDs.\n"
            "2. DeepSeek Structured JSON Output với ràng buộc trích dẫn Evidence ID.\n"
            "3. Citation Gate: Backend kiểm tra đối chiếu allow-list. Nếu DeepSeek bịa ra ID lạ -> Chặn ngay lập tức và kích hoạt 1 lượt Repair Request tự sửa sai. Bất biến: 100% câu hỏi và đánh giá đều truy nguyên được nguồn gốc!\""
        ),
        17: (
            "[07:00 - 08:00] SLIDE 17: VÒNG ĐỜI TÀI LIỆU RAG & HYBRID RETRIEVAL\n\n"
            "Lời nói:\n"
            "\"Kho tri thức 86 Rule Markdown (15 ngành, 60 profile STAR/BARS).\n"
            "Qdrant kết hợp song song:\n"
            "• Dense Embeddings: FastEmbed Multilingual MPNet.\n"
            "• Sparse Lexical: Thuật toán BM25 bắt từ khóa kỹ thuật.\n"
            "Hợp nhất bằng Reciprocal Rank Fusion (RRF), đảm bảo truy hồi tài liệu vừa chuẩn ngữ nghĩa vừa chính xác từ khóa.\""
        ),
        18: (
            "[08:00 - 09:30] SLIDE 18: PIPELINE TIẾNG NÓI ĐA TẦNG (REALTIME & HẬU KỲ)\n\n"
            "Lời nói:\n"
            "\"Kính thưa Thầy Cô, pipeline tiếng nói của InterV được thiết kế chuyên biệt theo 2 pha rõ rệt:\n\n"
            "1. PHA 1: LUỒNG REALTIME TRONG PHỎNG VẤN (TƯƠNG TÁC MƯỢT MÀ, KHÔNG RỜI RẠC)\n"
            "• Luồng âm thanh từ microphone được AudioWorklet thu và truyền trực tiếp qua WebSocket tới AssemblyAI để thực hiện STT theo thời gian thực.\n"
            "• Cơ chế Fallback cục bộ: Khi gặp sự cố mạng hoặc lỗi API, hệ thống tự động fallback về model Faster-Whisper chạy trực tiếp trên Python backend.\n"
            "• Mục đích: Luôn bảo đảm có dữ liệu transcript tức thời cho AI Provider (DeepSeek) tiếp tục hiểu ngữ cảnh, nói tiếp qua TTS và sinh các câu hỏi thích ứng (Lookahead Engine) liền mạch, không làm câu hỏi bị rời rạc hay ngắt quãng phiên phỏng vấn.\n\n"
            "2. PHA 2: LUỒNG XỬ LÝ HẬU KỲ SAU PHỎNG VẤN (ĐÁNH GIÁ CHUYÊN SÂU & COACHING)\n"
            "• Sau khi phỏng vấn xong, toàn bộ các đoạn audio gốc đã lưu sẽ được xử lý tuần tự (sequential batch) bằng Faster-Whisper để tái tạo STT chính xác cao nhất.\n"
            "• Hệ thống áp dụng các thuật toán xử lý tín hiệu âm thanh để đo lường khách quan nhịp độ, tốc độ nói (WPM), và các khoảng trống/khoảng lặng (pause duration) trong câu nói.\n"
            "• Tiếp theo, audio được đưa qua mô hình SenseVoice theo nguyên tắc 'Observation only' để bóc tách: LID (ngôn ngữ), SER (cảm xúc giọng nói), và AED (sự kiện âm thanh: tiếng cười, tiếng ho, thở dài, từ đệm/filler words).\n"
            "• Dữ liệu này được tổng hợp để đánh giá khách quan mức độ tự tin, tính lưu loát, và phong thái diễn đạt nhằm hỗ trợ coaching cho ứng viên mà không suy diễn tâm lý tùy tiện.\""
        ),
        19: (
            "[09:30 - 10:15] SLIDE 19: AN TOÀN, QUAN SÁT & KHẢ NĂNG PHỤC HỒI\n\n"
            "Lời nói:\n"
            "\"5 lớp bảo vệ: Auth/RBAC, Schema validation Zod/Pydantic, Timeout/Retry circuit breaker, Audit log/Usage tracking, Graceful fallback ở mọi khâu.\""
        ),
        20: (
            "[10:15 - 10:45] SLIDE 20: HỢP ĐỒNG TYPED gRPC\n\n"
            "Lời nói:\n"
            "\"Typed RPC 16 phương thức ràng buộc giữa TypeScript BFF và Python Backend. Dữ liệu chỉ được lưu trữ hoặc hiển thị sau khi vượt qua schema validation.\""
        ),
        21: (
            "[10:45 - 10:50] PHẦN 04: SẢN PHẨM & TRẢI NGHIỆM\n\n"
            "Lời nói:\n"
            "\"Tiếp theo là Phần 4: Sản phẩm và các luồng trải nghiệm thực tế trên hệ thống InterV.\""
        ),
        22: (
            "[10:50 - 11:10] SLIDE 22: MA TRẬN PHÂN QUYỀN RBAC\n\n"
            "Lời nói:\n"
            "\"Candidate, Recruiter, Admin. Ba vai trò phân quyền độc lập, AI tuyệt đối không có quyền tự động ra quyết định tuyển dụng hay loại bỏ ứng viên.\""
        ),
        23: (
            "[11:10 - 11:30] SLIDE 23: STATE MACHINE PHIÊN LUYỆN TẬP\n\n"
            "Lời nói:\n"
            "\"Vòng đời nghiêm ngặt: Draft -> Ready -> Recording -> Processing -> Reviewed -> Completed. Có cơ chế resume an toàn khi gặp sự cố mạng.\""
        ),
        24: (
            "[11:30 - 11:45] SLIDE 24: VÒNG ĐỜI TUYỂN DỤNG RECRUITER\n\n"
            "Lời nói:\n"
            "\"Job -> Campaign -> Invitation -> Interview -> Evidence -> Human Final Review. Recruiter sở hữu toàn quyền quyết định cuối cùng.\""
        ),
        25: (
            "[11:45 - 12:00] SLIDE 25: CƠ CHẾ FALLBACK TOÀN DIỆN\n\n"
            "Lời nói:\n"
            "\"Minh họa đường đi chính (Primary) và đường dự phòng (Fallback) tại mọi điểm trọng yếu: RAG fallback baseline, STT fallback Faster-Whisper, Audio fallback text-only.\""
        ),
        26: (
            "[12:00 - 12:20] SLIDE 26: STACK CÔNG NGHỆ\n\n"
            "Lời nói:\n"
            "\"Next.js 15, TypeScript, Tailwind v4, MongoDB ReplicaSet, Python gRPC (16 RPCs), DeepSeek, Qdrant RAG, AssemblyAI, Faster-Whisper, SenseVoice, Edge/Vbee TTS.\""
        ),
        27: (
            "[12:20 - 12:25] PHẦN 05: BẰNG CHỨNG & XÁC MINH\n\n"
            "Lời nói:\n"
            "\"Tiếp theo là Phần 5: Các bằng chứng kỹ thuật, quy mô hiện vật và kết quả kiểm thử toàn diện.\""
        ),
        28: (
            "[12:25 - 13:10] SLIDE 28: QUY MÔ HIỆN VẬT TRIỂN KHAI\n\n"
            "Lời nói:\n"
            "\"Con số thực tế từ repository:\n"
            "• Hơn 72.000 dòng mã nguồn logic.\n"
            "• 61 route files, 74 HTTP methods, 38 trang frontend.\n"
            "• 19 Mongoose models, 16 gRPC RPCs, 86 Rule files (15 ngành, 60 profile STAR/BARS).\""
        ),
        29: (
            "[13:10 - 13:40] SLIDE 29: XÁC MINH KỸ THUẬT ĐẠT 100%\n\n"
            "Lời nói:\n"
            "\"Toàn bộ 51/51 Backend tests đạt 100% trong 3.28 giây, 16/16 RPC contract tests pass, Next.js build đạt production-ready, zero-mock runtime.\""
        ),
        30: (
            "[13:40 - 13:55] SLIDE 30: BA ĐIỂM KIỂM SOÁT DEEPSEEK\n\n"
            "Lời nói:\n"
            "\"Kiểm soát Trước - Trong - Sau khi sinh: Pre-generation Context Grounding, In-generation JSON Schema, Post-generation Citation Gate & 1-shot Repair.\""
        ),
        31: (
            "[13:55 - 14:15] SLIDE 31: BẤT BIẾN AN TOÀN SENSEVOICE\n\n"
            "Lời nói:\n"
            "\"Observation-only (chỉ coaching, không chẩn đoán tâm lý), Neutral Default 50 nếu thiếu tín hiệu, không bao giờ tự ý trừ điểm của ứng viên.\""
        ),
        32: (
            "[14:15 - 14:20] PHẦN 06: KẾT LUẬN & ĐÓNG GÓP\n\n"
            "Lời nói:\n"
            "\"Cuối cùng là Phần 6: Đóng góp chính, giới hạn nghiên cứu và lộ trình kiểm soát của InterV.\""
        ),
        33: (
            "[14:20 - 14:35] SLIDE 33: BỐN ĐÓNG GÓP CHÍNH\n\n"
            "Lời nói:\n"
            "\"1. Kiến trúc Lookahead thời gian thực 0ms latency.\n"
            "2. Grounded Generation 3 cổng loại bỏ hoàn toàn ảo giác.\n"
            "3. Multimodal Speech Pipeline 2 pha an toàn, có fallback.\n"
            "4. Hệ thống Provenance & Rule Catalog 86 files chuẩn hóa theo STAR/BARS.\""
        ),
        34: (
            "[14:35 - 14:50] SLIDE 34: GIỚI HẠN NGHIÊN CỨU\n\n"
            "Lời nói:\n"
            "\"Đề tài trung thực xác nhận: Đạt tính nhất quán và vững chắc kỹ thuật, nhưng cần đánh giá Validity/Fairness thực nghiệm trên tập mẫu lớn trước khi triển khai tuyển dụng diện rộng.\""
        ),
        35: (
            "[14:50 - 14:55] SLIDE 35: BỐN CỔNG KIỂM SOÁT TRIỂN KHAI\n\n"
            "Lời nói:\n"
            "\"Lộ trình mở rộng qua 4 gates: Practice Playground -> Shadow Scoring -> Human-overridden Pilot -> Restricted Production.\""
        ),
        36: (
            "[14:55 - 14:58] SLIDE 36: KẾT LUẬN\n\n"
            "Lời nói:\n"
            "\"InterV đã đạt nền tảng kỹ thuật khả thi, an toàn và sẵn sàng phục vụ luyện tập và hỗ trợ tuyển dụng có kiểm soát.\""
        ),
        37: (
            "[14:58 - 15:00] SLIDE 37: CẢM ƠN & TIẾP NHẬN PHẢN BIỆN\n\n"
            "Lời nói:\n"
            "\"Em xin trân trọng cảm ơn Thầy ThS. Đặng Văn Lực và quý Thầy Cô trong Hội đồng đã chú ý theo dõi. Em xin phép được tiếp nhận câu hỏi phản biện!\""
        )
    }

    for idx in range(len(prs.slides)):
        slide_num = idx + 1
        slide = prs.slides[idx]
        note_text = speaking_notes_37.get(slide_num, "")
        notes_slide = slide.notes_slide
        if notes_slide and notes_slide.notes_text_frame:
            notes_slide.notes_text_frame.text = note_text

    # Save to v20 and v18 (with graceful handling if open)
    v20_path = r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v20.pptx"
    v18_path = r"d:\project\InterV\Review\KLTN_InterV_LeMinhDuy_v18.pptx"
    
    for path in [v20_path, v18_path]:
        try:
            prs.save(path)
            print(f"Saved {os.path.basename(path)} successfully.")
        except PermissionError:
            alt_path = path.replace(".pptx", "_new.pptx")
            prs.save(alt_path)
            print(f"Notice: {os.path.basename(path)} is locked by PowerPoint. Saved to {os.path.basename(alt_path)} instead.")

    return prs

def update_markdown_and_scripts_37(prs):
    # 1. Update current_slides_analysis.md
    md_path = r"d:\project\InterV\Review\current_slides_analysis.md"
    new_md_lines = [f"# PHÂN TÍCH VÀ NỘI DUNG TỪNG SLIDE (KLTN_InterV_LeMinhDuy_v20.pptx - {len(prs.slides)} Slides)\n\n"]
    for idx, slide in enumerate(prs.slides):
        new_md_lines.append(f"## Slide {idx + 1}\n\n")
        for s in slide.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    txt = p.text.strip()
                    if txt:
                        new_md_lines.append(f"- {txt}\n")
        new_md_lines.append("\n---\n\n")
        
    with open(md_path, "w", encoding="utf-8") as f:
        f.writelines(new_md_lines)
    print(f"Successfully updated {md_path}")
    
    # 2. Update slide_notes_dump.json
    json_path = r"d:\project\InterV\Review\slide_notes_dump.json"
    notes_list = []
    for idx, slide in enumerate(prs.slides):
        note = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide and slide.notes_slide.notes_text_frame else ""
        notes_list.append({
            "slide_num": idx + 1,
            "notes": note
        })
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(notes_list, f, ensure_ascii=False, indent=2)
    print(f"Successfully updated {json_path}")

if __name__ == "__main__":
    prs = build_master_deck()
    update_markdown_and_scripts_37(prs)
    print("ALL 37 SLIDES READY!")
